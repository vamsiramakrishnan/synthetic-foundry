from __future__ import annotations

from datetime import datetime, timezone

from worldloom.artifact_shape import plan_artifact_shape
from worldloom.artifact_shape_materialize import materialize_artifact_shape
from worldloom.eval_design import ArtifactShapeRequirement
from worldloom.evidence_locator import DocxLocator, PptxLocator, XlsxLocator
from worldloom.models import Authority, CanonicalFact, Quantity


def _fact() -> CanonicalFact:
    return CanonicalFact(
        id="F-EVIDENCE",
        kind="q3_spend",
        subject="finance:apac",
        value=Quantity(amount=881000, unit="SGD"),
        valid_from=datetime(2026, 8, 31, tzinfo=timezone.utc),
        authority=Authority.SYSTEM_OF_RECORD,
        source_system="erp",
    )


def test_pptx_shape_has_exact_slide_notes_and_real_file_weight(tmp_path) -> None:  # type: ignore[no-untyped-def]
    from pptx import Presentation

    fact = _fact()
    requirement = ArtifactShapeRequirement(
        artifact_type="pptx",
        slides=30,
        versions=3,
        speaker_note_slides=12,
        hidden_slides=2,
        native_charts=5,
        file_size_bytes=300_000,
        evidence_index=27,
        evidence_modality="speaker_notes",
    )
    plan = plan_artifact_shape(
        requirement,
        artifact_id="deck:q3-review",
        fact_ids=(fact.id,),
    )
    path = tmp_path / "large.pptx"

    result = materialize_artifact_shape(
        plan,
        path,
        title="Q3 Business Review",
        facts={fact.id: fact},
    )
    presentation = Presentation(path)
    locator = result.evidence.locator if result.evidence else None

    assert result.slides == 30
    assert len(presentation.slides) == 30
    assert result.file_size_bytes >= 300_000
    assert isinstance(locator, PptxLocator)
    assert locator.slide == 27
    assert "881000" in presentation.slides[26].notes_slide.notes_text_frame.text
    assert sum(1 for slide in presentation.slides if slide._element.get("show") == "0") >= 2  # noqa: SLF001
    assert sum(
        1
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_chart", False)
    ) >= 5


def test_docx_shape_uses_explicit_page_breaks_and_exact_paragraph_locator(tmp_path) -> None:  # type: ignore[no-untyped-def]
    from docx import Document

    fact = _fact()
    requirement = ArtifactShapeRequirement(
        artifact_type="docx",
        pages=12,
        paragraphs=120,
        versions=5,
        file_size_bytes=180_000,
        evidence_index=9,
        evidence_modality="text",
    )
    plan = plan_artifact_shape(
        requirement,
        artifact_id="doc:operating-review",
        fact_ids=(fact.id,),
    )
    path = tmp_path / "long.docx"

    result = materialize_artifact_shape(
        plan,
        path,
        title="Operating Review",
        facts={fact.id: fact},
    )
    document = Document(path)
    locator = result.evidence.locator if result.evidence else None

    assert result.pages == 12
    assert result.paragraphs == 120
    assert result.file_size_bytes >= 180_000
    assert isinstance(locator, DocxLocator)
    assert locator.approximate_page == 9
    assert "881000" in document.paragraphs[locator.paragraph - 1].text


def test_xlsx_shape_is_logically_large_sparse_and_formula_gradeable(tmp_path) -> None:  # type: ignore[no-untyped-def]
    from openpyxl import load_workbook

    fact = _fact()
    requirement = ArtifactShapeRequirement(
        artifact_type="xlsx",
        sheets=4,
        rows_per_sheet=2_000,
        columns_per_sheet=40,
        formulas=1_000,
        comments=20,
        file_size_bytes=400_000,
        evidence_index=1_750,
        evidence_modality="formula",
    )
    plan = plan_artifact_shape(
        requirement,
        artifact_id="book:q3-finance",
        fact_ids=(fact.id,),
    )
    path = tmp_path / "huge.xlsx"

    result = materialize_artifact_shape(
        plan,
        path,
        title="Q3 Finance Workbook",
        facts={fact.id: fact},
    )
    workbook = load_workbook(path, read_only=False, data_only=False)
    locator = result.evidence.locator if result.evidence else None

    assert result.sheets == 4
    assert result.rows_per_sheet == 2_000
    assert result.columns_per_sheet == 40
    assert result.logical_cells == 320_000
    assert result.formulas == 1_000
    assert result.comments == 20
    assert result.file_size_bytes >= 400_000
    assert isinstance(locator, XlsxLocator)
    assert locator.cell == "B1750"
    assert workbook[locator.sheet][locator.cell].value == "=A1750"
    assert workbook[locator.sheet]["A1750"].value == 881000
    assert workbook["Sheet4"].max_row == 2_000
    assert workbook["Sheet4"].max_column == 40
