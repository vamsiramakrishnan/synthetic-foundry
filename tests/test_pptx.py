"""Step 5.4: the PPTX vertical slice.

Per ``docs/artifact-compiler.md`` section 15, the PPTX Gate is five concrete
claims, and this file's job is to make each one an assertion rather than an
impression:

* one executive deck renders from existing IR
* no facts or values are introduced by the renderer
* deck content agrees with Markdown, DOCX, and XLSX
* the same seed and ledger produce byte-identical normalised OOXML
* all shapes remain in bounds and non-overlapping

The bounds/overlap/font-floor checks are written once, as reusable helpers,
and run over every slide of every deck built in this file — the gate is a
property of the deck as a whole, not of one slide chosen to look good.
"""

from __future__ import annotations

import io
import sys
import zipfile
from datetime import datetime

import pptx as python_pptx
import pytest

from worldloom import MonthEndClose, RetailWorld, World
from worldloom.models import Cell, Chart, ChartKind, Column, Row, Table
from worldloom.narrative import DeterministicProvider, references
from worldloom.render import RenderError, ooxml
from worldloom.render import docx as docx_renderer
from worldloom.render import pptx as pptx_renderer
from worldloom.render.values import format_value

PERIOD = "2026-03"


@pytest.fixture(scope="module")
def rendered() -> World:
    world = RetailWorld(seed=8128).build().run(
        MonthEndClose(period=PERIOD, include_operational_incident=True)
    )
    return world.narrate(DeterministicProvider()).render("pptx", "markdown", "docx", "xlsx")


def _files(world: World, suffix: str) -> dict[str, bytes]:
    return {r.artifact_id: r.payload for r in world._rendered if r.path.endswith(suffix)}


def _deck(payload: bytes):  # type: ignore[no-untyped-def]
    return python_pptx.Presentation(io.BytesIO(payload))


def _executive_summary_ir(rendered: World):  # type: ignore[no-untyped-def]
    return next(
        ir for ir in rendered.artifact_irs
        if rendered.artifact_intents.by_id(ir.intent_id).artifact_type == "executive_summary"
    )


# ---------------------------------------------------------------------------
# Reusable structural checks — run over every slide of every deck below
# ---------------------------------------------------------------------------


def _shape_rects(slide) -> list[tuple[int, int, int, int, int]]:  # type: ignore[no-untyped-def]
    return [
        (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height, shape.shape_id)
        for shape in slide.shapes
    ]


def assert_shapes_in_bounds_and_non_overlapping(prs) -> None:  # type: ignore[no-untyped-def]
    """The PPTX Gate's structural criterion, literally: every shape's declared
    box sits inside the canvas, and no two shapes on the same slide overlap.

    A deterministic check on declared geometry, not a rendered screenshot —
    which is exactly what 14.A's "structural validation" step (as opposed to
    the optional raster one) means.
    """
    width, height = prs.slide_width, prs.slide_height
    for index, slide in enumerate(prs.slides):
        rects = _shape_rects(slide)
        for left, top, right, bottom, shape_id in rects:
            assert left >= 0 and top >= 0, f"slide {index} shape {shape_id} starts out of bounds"
            assert right <= width and bottom <= height, (
                f"slide {index} shape {shape_id} exceeds the {width}x{height} canvas: "
                f"right={right} bottom={bottom}"
            )
        for i in range(len(rects)):
            for j in range(i + 1, len(rects)):
                l1, t1, r1, b1, id1 = rects[i]
                l2, t2, r2, b2, id2 = rects[j]
                overlaps = l1 < r2 and l2 < r1 and t1 < b2 and t2 < b1
                assert not overlaps, f"slide {index}: shapes {id1} and {id2} overlap"


def _text_frames(shape):  # type: ignore[no-untyped-def]
    """Every text frame a shape carries — its own, or one per table cell."""
    if shape.has_table:
        return [cell.text_frame for row in shape.table.rows for cell in row.cells]
    if shape.has_text_frame:
        return [shape.text_frame]
    return []


def assert_font_floor(prs, floor_pt: int = pptx_renderer.MIN_FONT_PT) -> None:  # type: ignore[no-untyped-def]
    """No run anywhere in the deck is smaller than *floor_pt*.

    Every run this renderer writes carries an explicit size (`_write` and
    `_style_cell` never leave it to theme inheritance — see their docstrings),
    so a missing size here is itself a defect worth failing on, not something
    to skip past.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            for frame in _text_frames(shape):
                for paragraph in frame.paragraphs:
                    for run in paragraph.runs:
                        assert run.font.size is not None, f"a run has no explicit size: {run.text!r}"
                        assert run.font.size.pt >= floor_pt, (
                            f"{run.text!r} is {run.font.size.pt}pt, below the {floor_pt}pt floor"
                        )


def _deck_text(prs) -> str:  # type: ignore[no-untyped-def]
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            for frame in _text_frames(shape):
                if frame.text:
                    parts.append(frame.text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Gate: one executive deck renders from existing IR
# ---------------------------------------------------------------------------


def test_the_executive_summary_renders_to_pptx(rendered: World) -> None:
    decks = _files(rendered, ".pptx")
    ir = _executive_summary_ir(rendered)
    assert ir.id in decks, "the retail-close executive summary should have a deck"
    assert len(decks) == 1, "this vertical slice targets only executive_summary — see HANDLES"

    deck = _deck(decks[ir.id])
    assert len(deck.slides) >= 6, "title, three beats, an appendix, and a close is at least six slides"


def test_only_executive_summary_is_handled(rendered: World) -> None:
    assert pptx_renderer.HANDLES == frozenset({"executive_summary"})
    produced_types = {
        rendered.artifact_intents.by_id(
            next(ir.intent_id for ir in rendered.artifact_irs if ir.id == artifact_id)
        ).artifact_type
        for artifact_id in _files(rendered, ".pptx")
    }
    assert produced_types == {"executive_summary"}


def test_one_artifact_has_one_basename_across_formats(rendered: World) -> None:
    deck = {r.artifact_id: r.path for r in rendered._rendered if r.path.endswith(".pptx")}
    markdown_files = {r.artifact_id: r.path for r in rendered._rendered if r.path.endswith(".md")}
    for artifact_id, path in deck.items():
        assert markdown_files[artifact_id].removesuffix(".md") == path.removesuffix(".pptx")


def test_slide_count_matches_the_plan(rendered: World) -> None:
    """A round trip against the renderer's own plan, not a hand counted number.

    Pagination means one `SlidePlan` can become several physical slides — the
    plan itself only records that a section carries a 12-row table, not how
    many slides that becomes, so this recomputes the same chunking the
    renderer does rather than hard-coding today's row count.
    """
    ir = _executive_summary_ir(rendered)
    plan = pptx_renderer._plan(ir)

    expected = 0
    for slide_plan in plan.slides:
        if slide_plan.table is None:
            expected += 1
            continue
        rows = len(slide_plan.table.rows)
        columns = len(slide_plan.table.columns)
        if slide_plan.component_id == "finance.metric_strip" and 1 <= rows <= 6 and columns <= 3:
            expected += 1
        else:
            cap = pptx_renderer._MAX_TABLE_ROWS_PER_SLIDE
            expected += max(1, -(-rows // cap))  # ceiling division

    deck = _deck(_files(rendered, ".pptx")[ir.id])
    assert len(deck.slides) == expected


def test_the_deck_uses_the_compiler(rendered: World) -> None:
    """The whole point of `plan_from_ir`/`compose` existing: the component ids
    a real deck is built from are the compiler's choice, not a hard-coded
    outline this module invented on its own."""
    from worldloom.compiler.components import REGISTRY

    pptx_component_ids = {spec.component_id for spec in REGISTRY if "pptx" in spec.supported_formats}

    ir = _executive_summary_ir(rendered)
    plan = pptx_renderer._plan(ir)
    content_components = {s.component_id for s in plan.slides if s.kind == "content"}
    assert content_components, "the plan should carry at least one compiler-assigned component"
    assert content_components <= pptx_component_ids


def test_structural_gate_holds_for_the_real_deck(rendered: World) -> None:
    ir = _executive_summary_ir(rendered)
    deck = _deck(_files(rendered, ".pptx")[ir.id])
    assert_shapes_in_bounds_and_non_overlapping(deck)
    assert_font_floor(deck)


# ---------------------------------------------------------------------------
# Gate: no facts or values are introduced by the renderer
# ---------------------------------------------------------------------------


def _expected_haystack(ir, facts) -> str:  # type: ignore[no-untyped-def]
    """Every string this artifact's IR actually contains, concatenated.

    Anything with a digit that shows up on a slide has to be a substring of
    this — if it is not, the renderer typed a figure nobody gave it, which is
    exactly what the bare-number rule in `AGENTS.md` exists to catch, applied
    to the one format that rule has never been checked against before.
    """
    parts = [ir.title, ir.subtitle or ""]
    parts.extend(str(v) for v in ir.metadata.values())
    for section in ir.sections:
        parts.append(section.heading)
        if section.body:
            parts.append(references.substitute(section.body, facts))
        if section.table is not None:
            parts.append(section.table.title)
            if section.table.note:
                parts.append(section.table.note)
            for row in section.table.rows:
                parts.append(row.label)
                for column in section.table.columns:
                    cell = row.cells.get(column.key)
                    if cell is None:
                        continue
                    parts.append(format_value(cell.value, column.number_format))
                    if cell.fact_id:
                        parts.append(cell.fact_id)
    return "\n".join(parts)


def test_no_facts_or_values_are_introduced(rendered: World) -> None:
    facts = {fact.id: fact for fact in rendered.facts}
    ir = _executive_summary_ir(rendered)
    deck = _deck(_files(rendered, ".pptx")[ir.id])

    haystack = _expected_haystack(ir, facts)
    checked = 0
    for slide in deck.slides:
        for shape in slide.shapes:
            for frame in _text_frames(shape):
                text = frame.text
                if any(ch.isdigit() for ch in text):
                    assert text in haystack, f"digit-bearing text not traceable to the IR: {text!r}"
                    checked += 1
    assert checked > 5, f"only checked {checked} numeric fragments — the sweep is not covering the deck"


def test_no_fact_reference_survives_into_the_deck(rendered: World) -> None:
    for payload in _files(rendered, ".pptx").values():
        text = _deck_text(_deck(payload))
        assert "{{fact:" not in text
        assert "[missing " not in text


def test_prose_still_contains_no_bare_number(rendered: World) -> None:
    """The same invariant `test_docx.py` checks: figures arrive by
    substitution, never typed into the section body itself."""
    facts = {fact.id: fact for fact in rendered.facts}
    ir = _executive_summary_ir(rendered)
    for section in ir.sections:
        if section.body:
            assert not references.bare_numbers(section.body), (
                f"{ir.id}/{section.heading} restated a figure"
            )
            assert references.substitute(section.body, facts) != section.body


# ---------------------------------------------------------------------------
# Gate: deck content agrees with Markdown, DOCX, and XLSX
# ---------------------------------------------------------------------------


def test_deck_and_markdown_carry_the_same_prose(rendered: World) -> None:
    facts = {fact.id: fact for fact in rendered.facts}
    ir = _executive_summary_ir(rendered)
    deck_text = _deck_text(_deck(_files(rendered, ".pptx")[ir.id]))

    checked = 0
    for section in ir.sections:
        assert section.heading in deck_text, f"missing section heading {section.heading!r}"
        if section.body:
            resolved = references.substitute(section.body, facts)
            for sentence in resolved.split(". "):
                fragment = sentence.strip().rstrip(".")
                if len(fragment) > 15:
                    assert fragment in deck_text, f"prose dropped from the deck: {fragment!r}"
                    checked += 1
    assert checked > 0, "no prose fragment was long enough to compare"


def test_deck_and_docx_agree_on_every_table_cell(rendered: World) -> None:
    """Both are projections of the same IR — the renderer that reads a table
    and the renderer that reads the same table cannot show different figures."""
    ir = _executive_summary_ir(rendered)
    deck_text = _deck_text(_deck(_files(rendered, ".pptx")[ir.id]))
    docx_text = docx_renderer.render(ir, {f.id: f for f in rendered.facts})
    import docx as python_docx

    document = python_docx.Document(io.BytesIO(docx_text))
    docx_body = "\n".join(p.text for p in document.paragraphs) + "\n".join(
        cell.text for table in document.tables for row in table.rows for cell in row.cells
    )

    checked = 0
    for section in ir.sections:
        if section.table is None:
            continue
        for row in section.table.rows:
            for column in section.table.columns:
                cell = row.cells.get(column.key)
                if cell is None or cell.value is None:
                    continue
                text = format_value(cell.value, column.number_format)
                assert text in deck_text, f"cell {row.key}/{column.key} = {text!r} missing from the deck"
                assert text in docx_body, f"cell {row.key}/{column.key} = {text!r} missing from the docx twin"
                checked += 1
    assert checked > 20, f"only checked {checked} cells"


def test_deck_headline_agrees_with_the_fact_ledger(rendered: World) -> None:
    """The ledger is what the XLSX workbook's formulas evaluate to
    (`test_render.py::test_markdown_and_xlsx_report_the_same_numbers` proves
    that equality for the workbook itself) — so agreeing with the ledger here
    is agreeing with XLSX by the transitive property both share one IR."""
    ir = _executive_summary_ir(rendered)
    deck_text = _deck_text(_deck(_files(rendered, ".pptx")[ir.id]))

    revenue = rendered.facts.where(kind="financial.revenue.actual", subject=rendered.company.id).one()
    rendered_value = references.render_value(revenue)
    assert rendered_value in deck_text


# ---------------------------------------------------------------------------
# Gate: the same seed and ledger produce byte-identical normalised OOXML
# ---------------------------------------------------------------------------


def test_rendering_twice_is_byte_identical(rendered: World) -> None:
    """Two separate `render()` calls, not one render compared to itself, and
    not two builds a second apart — see the module note on why that would
    pass by luck. Both calls happen in this process, well within one second,
    which is exactly the case that would previously have hidden a clock
    leaking into the archive."""
    ir = _executive_summary_ir(rendered)
    facts = {fact.id: fact for fact in rendered.facts}
    assert pptx_renderer.render(ir, facts) == pptx_renderer.render(ir, facts)


def test_two_independent_builds_are_byte_identical(rendered: World) -> None:
    """The stronger claim: not the same IR object rendered twice, but two
    separate worlds built from the same seed, narrated, and rendered."""

    def build_pptx() -> bytes:
        world = RetailWorld(seed=8128).build().run(
            MonthEndClose(period=PERIOD, include_operational_incident=True)
        )
        world = world.narrate(DeterministicProvider()).render("pptx")
        ir = next(
            ir for ir in world.artifact_irs
            if world.artifact_intents.by_id(ir.intent_id).artifact_type == "executive_summary"
        )
        return next(r.payload for r in world._rendered if r.artifact_id == ir.id)

    assert build_pptx() == build_pptx()


def test_no_clock_reaches_the_deck(rendered: World) -> None:
    """Checked as an invariant, not by rendering twice — the same reasoning
    `test_docx.py` gives: two renders landing in the same second would pass
    even with a live clock, which is exactly how this defect hid in XLSX."""
    for artifact_id, payload in _files(rendered, ".pptx").items():
        with zipfile.ZipFile(io.BytesIO(payload)) as archive:
            stamps = {info.date_time for info in archive.infolist()}
        assert stamps == {ooxml.EPOCH}, f"{artifact_id}: archive carries wall-clock entries {stamps}"

        ir = next(ir for ir in rendered.artifact_irs if ir.id == artifact_id)
        properties = _deck(payload).core_properties
        expected = datetime.fromisoformat(ir.metadata["worldloom_created"])
        # python-pptx's own W3CDTF reader drops the trailing "Z" rather than
        # treating it as UTC (`pptx.oxml.coreprops._parse_W3CDTF_to_datetime`
        # only special-cases a numeric `+hh:mm` offset), so it hands back a
        # naive datetime where python-docx hands back a UTC-aware one — a
        # library difference in the *reader*, not in the timestamp this
        # renderer wrote, which is why the XML assertion above compares the
        # raw string and this one strips tzinfo before comparing.
        assert properties.created == expected.replace(tzinfo=None)
        assert properties.modified == expected.replace(tzinfo=None)
        assert properties.created.year == 2026, "python-pptx's template date leaked through"


# ---------------------------------------------------------------------------
# Dependency handling
# ---------------------------------------------------------------------------


def test_missing_python_pptx_raises_actionably(monkeypatch: pytest.MonkeyPatch) -> None:
    """Matches `docx.py`'s `_require_docx`: an actionable `RenderError`
    naming the install command, not a bare `ImportError` from deep inside the
    renderer."""
    monkeypatch.setitem(sys.modules, "pptx", None)
    with pytest.raises(RenderError, match="pip install"):
        pptx_renderer._require_pptx()


# ---------------------------------------------------------------------------
# Native charts — the retail-close executive summary declares none today
# (see the task's own finding: an executive deck with zero charts), so these
# exercise the general machinery every `Chart` kind shares directly, the same
# way the "Components exercised directly" tests below reach a component the
# real deck happens not to select.
# ---------------------------------------------------------------------------

def _chart_table_and_chart(kind: ChartKind, *, by_row: bool = False, rows: list[str] | None = None):  # type: ignore[no-untyped-def]
    table = Table(
        key="pnl", title="P&L", columns=[
            Column(key="budget", label="Budget", number_format="#,##0;(#,##0)"),
            Column(key="actual", label="Actual", number_format="#,##0;(#,##0)"),
        ],
        rows=[
            Row(key="food", label="Food", cells={"budget": Cell(value=100.0), "actual": Cell(value=110.0)}),
            Row(key="apparel", label="Apparel", cells={"budget": Cell(value=200.0), "actual": Cell(value=190.0)}),
            Row(key="group", label="Group", emphasis=True,
                cells={"budget": Cell(value=300.0), "actual": Cell(value=300.0)}),
        ],
    )
    chart = Chart(
        key="test", title="Test chart", kind=kind, table="pnl",
        series=["budget", "actual"], rows=rows if rows is not None else [], by_row=by_row,
        category_axis="Division", value_axis="AUD thousands", note="A note on the chart.",
    )
    return table, chart


def _draw_chart_deck(kind: ChartKind, **kwargs):  # type: ignore[no-untyped-def]
    table, chart = _chart_table_and_chart(kind, **kwargs)
    presentation = _blank_deck()
    pptx_renderer._draw_chart(presentation, chart, table, footer_text="Test · Deck")
    return presentation, table, chart


@pytest.mark.parametrize("kind", [ChartKind.COLUMN, ChartKind.BAR, ChartKind.LINE, ChartKind.PIE])
def test_every_chart_kind_draws_in_bounds_with_the_font_floor(kind: ChartKind) -> None:
    presentation, _, _ = _draw_chart_deck(kind)
    assert len(presentation.slides) == 1
    assert_shapes_in_bounds_and_non_overlapping(presentation)
    assert_font_floor(presentation)


def test_chart_series_resolves_categories_and_values_off_the_table() -> None:
    """`Chart` never introduces a number — `_chart_series` reads nothing but
    cells already in the table it is handed."""
    table, chart = _chart_table_and_chart(ChartKind.COLUMN)
    categories, series = pptx_renderer._chart_series(chart, table)
    assert categories == ["Food", "Apparel"], "the subtotal row should not be a category"
    assert dict((label, values) for label, values, _ in series) == {
        "Budget": [100.0, 200.0],
        "Actual": [110.0, 190.0],
    }


def test_empty_rows_means_every_row_that_is_not_a_subtotal() -> None:
    """`Chart.rows`'s own docstring: a chart that included the subtotal would
    double every bar."""
    table, chart = _chart_table_and_chart(ChartKind.COLUMN, rows=[])
    categories, _ = pptx_renderer._chart_series(chart, table)
    assert categories == ["Food", "Apparel"]
    assert "Group" not in categories


def test_by_row_reads_rows_as_series_and_columns_as_categories() -> None:
    """`Chart.by_row`'s own docstring: read the wrong way round, this renders
    without complaint as one point per series instead of one line per row."""
    table, chart = _chart_table_and_chart(ChartKind.LINE, by_row=True, rows=["food", "apparel"])
    categories, series = pptx_renderer._chart_series(chart, table)
    assert categories == ["Budget", "Actual"]
    assert [label for label, _, _ in series] == ["Food", "Apparel"]


def test_a_percentage_column_is_scaled_the_same_way_the_workbook_scales_it() -> None:
    """A percentage fact stored as 24.94 should plot as 0.2494 — the same
    conversion `xlsx.py::render` applies before Excel's percent format
    multiplies it back out, so an axis labelled "0.00%" reads 24.94%, not
    2494%."""
    table = Table(
        key="margins", title="Margins",
        columns=[Column(key="gm_pct", label="GM%", number_format="0.00%")],
        rows=[Row(key="food", label="Food", cells={"gm_pct": Cell(value=24.94)})],
    )
    chart = Chart(key="margin", title="Margin", kind=ChartKind.COLUMN, table="margins", series=["gm_pct"])
    _, series = pptx_renderer._chart_series(chart, table)
    assert series == [("GM%", [0.2494], "0.00%")]


def test_a_pie_chart_plots_only_its_first_series() -> None:
    """`ChartKind.PIE`'s own docstring: composition of a single total, one
    series only — a second series would silently double the wedges."""
    presentation, _, _ = _draw_chart_deck(ChartKind.PIE)
    slide = presentation.slides[0]
    charts = [shape.chart for shape in slide.shapes if shape.has_chart]
    assert len(charts) == 1
    assert len(charts[0].series) == 1


def test_a_chart_reading_a_different_table_than_its_section_is_skipped() -> None:
    """`Chart`'s own docstring: every value it plots is a cell already in the
    table beside it — a mismatched `chart.table` is dropped, not guessed at."""
    table, chart = _chart_table_and_chart(ChartKind.COLUMN)
    mismatched = chart.model_copy(update={"table": "not_this_table"})
    slide_plan = pptx_renderer.SlidePlan(
        kind="content", component_id="finance.variance_table", heading="Test", table=table,
        charts=(mismatched,),
    )
    presentation = _blank_deck()
    plan = pptx_renderer.PresentationPlan(title="T", subtitle=None, metadata={}, slides=(slide_plan,))
    pptx_renderer._draw_content(presentation, plan, slide_plan, {})
    # Only the table slide — no extra chart slide was added.
    assert len(presentation.slides) == 1


def test_a_chart_slide_is_added_after_the_table_when_it_matches() -> None:
    """Same dispatch, exercised through `_draw_content` rather than by
    calling `_draw_chart` directly — the table slide, then one chart slide."""
    table, chart = _chart_table_and_chart(ChartKind.COLUMN)
    slide_plan = pptx_renderer.SlidePlan(
        kind="content", component_id="finance.variance_table", heading="Divisional", table=table,
        charts=(chart,),
    )
    presentation = _blank_deck()
    plan = pptx_renderer.PresentationPlan(title="T", subtitle=None, metadata={}, slides=(slide_plan,))
    pptx_renderer._draw_content(presentation, plan, slide_plan, {})
    assert len(presentation.slides) == 2
    assert any(shape.has_chart for shape in presentation.slides[1].shapes)
    assert_shapes_in_bounds_and_non_overlapping(presentation)
    assert_font_floor(presentation)


def test_a_chart_deck_renders_twice_byte_identical() -> None:
    table, chart = _chart_table_and_chart(ChartKind.COLUMN)
    slide_plan = pptx_renderer.SlidePlan(
        kind="content", component_id="finance.variance_table", heading="Divisional", table=table,
        charts=(chart,),
    )
    plan = pptx_renderer.PresentationPlan(title="T", subtitle=None, metadata={}, slides=(slide_plan,))

    def build() -> bytes:
        presentation = _blank_deck()
        pptx_renderer._draw_content(presentation, plan, slide_plan, {})
        buffer = io.BytesIO()
        presentation.save(buffer)
        # A fixed stamp, not an omitted one. Without `created`, `normalise`
        # fixes archive entry dates and leaves every XML timestamp alone —
        # deliberately, and pinned by `test_ooxml`. A chart embeds a
        # data-source workbook that `xlsxwriter` stamps with
        # `datetime.now(timezone.utc)`, so an un-stamped normalise still
        # leaves a live clock two containers down: measured, two renders
        # 2.2s apart differed by one byte inside
        # `ppt/embeddings/Microsoft_Excel_Sheet1.xlsx`. This test therefore
        # passed or failed on whether both saves landed in the same second,
        # which is why it only ever failed on a loaded machine.
        return ooxml.normalise(buffer.getvalue(), created="2026-04-01T07:00:00+00:00")

    assert build() == build()


def test_no_clock_reaches_a_chart_deck_including_its_embedded_workbook() -> None:
    """The nested-clock defect `ooxml.normalise`'s own docstring describes:
    `python-pptx`'s chart API always embeds a small `.xlsx` workbook of the
    plotted values, and `xlsxwriter` stamps *that* workbook's own
    `docProps/core.xml` with `datetime.now()` — a second clock the top-level
    substitution alone never reaches."""
    table, chart = _chart_table_and_chart(ChartKind.COLUMN)
    slide_plan = pptx_renderer.SlidePlan(
        kind="content", component_id="finance.variance_table", heading="Divisional", table=table,
        charts=(chart,),
    )
    plan = pptx_renderer.PresentationPlan(title="T", subtitle=None, metadata={}, slides=(slide_plan,))
    presentation = _blank_deck()
    pptx_renderer._draw_content(presentation, plan, slide_plan, {})
    buffer = io.BytesIO()
    presentation.save(buffer)
    normalised = ooxml.normalise(buffer.getvalue())

    with zipfile.ZipFile(io.BytesIO(normalised)) as archive:
        stamps = {info.date_time for info in archive.infolist()}
        assert stamps == {ooxml.EPOCH}
        embedded = [n for n in archive.namelist() if n.endswith(".xlsx")]
        assert embedded, "a native pptx chart should carry its own embedded workbook"
        for name in embedded:
            with zipfile.ZipFile(io.BytesIO(archive.read(name))) as inner:
                inner_stamps = {info.date_time for info in inner.infolist()}
                assert inner_stamps == {ooxml.EPOCH}, f"{name}: embedded workbook still carries a wall clock"


# ---------------------------------------------------------------------------
# Components exercised directly — not every component the registry declares
# reaches the retail-close executive summary (see `draw_metric_cards`'s own
# docstring on why), so these test the shapes those components draw without
# waiting for an artifact type that happens to select them.
# ---------------------------------------------------------------------------


def _blank_deck():  # type: ignore[no-untyped-def]
    presentation = python_pptx.Presentation()
    presentation.slide_width = pptx_renderer.SLIDE_W
    presentation.slide_height = pptx_renderer.SLIDE_H
    return presentation


def test_metric_strip_component_is_in_bounds_and_non_overlapping() -> None:
    table = Table(
        key="headline",
        title="Headline measures",
        columns=[Column(key="value", label="Value", number_format="#,##0")],
        rows=[
            Row(key="revenue", label="Revenue", cells={"value": _cell(617_200)}),
            Row(key="gross_profit", label="Gross profit", cells={"value": _cell(151_325)}),
            Row(key="margin", label="Margin", cells={"value": _cell(24.52, fmt="0.00%")}),
        ],
    )
    presentation = _blank_deck()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    pptx_renderer.draw_metric_cards(slide, table)

    assert_shapes_in_bounds_and_non_overlapping(presentation)
    assert_font_floor(presentation)
    text = _deck_text(presentation)
    for row in table.rows:
        assert row.label in text
    assert "617,200" in text


def _cell(value, fmt: str | None = "#,##0"):  # type: ignore[no-untyped-def]
    from worldloom.models import Cell

    return Cell(value=value)


def test_decision_panel_and_schedule_prose_styles_are_in_bounds() -> None:
    """The accent-bar treatment `core.schedule`/`mgmt.decision_panel` use
    (`_PROSE_STYLE`) adds a second shape next to the body text — the one
    prose layout in this renderer where non-overlap is not simply "one box
    fills BODY", so it is worth checking on its own rather than only ever
    incidentally through the real deck."""
    presentation = _blank_deck()
    for component_id in ("core.schedule", "mgmt.decision_panel", "core.narrative", "core.position"):
        slide_plan = pptx_renderer.SlidePlan(
            kind="content",
            component_id=component_id,
            heading="Test section",
            body="Body copy for a synthetic section used only to exercise this component's geometry.",
        )
        pptx_renderer._draw_prose_content(presentation, slide_plan, {}, footer_text="Test · Deck")

    assert_shapes_in_bounds_and_non_overlapping(presentation)
    assert_font_floor(presentation)


def test_table_pagination_stays_in_bounds_for_a_wide_row_count() -> None:
    """A table well past one slide's row cap — the shape this renderer's own
    pagination exists for, exercised directly rather than only via whatever
    row count `retail-close` happens to produce today."""
    columns = [Column(key="value", label="Value", number_format="#,##0")]
    rows = [Row(key=f"r{i}", label=f"Row {i}", cells={"value": _cell(i * 1000)}) for i in range(23)]
    table = Table(key="wide", title="Wide table", columns=columns, rows=rows, note="A note on the data.")

    presentation = _blank_deck()
    pptx_renderer._draw_table_content(
        presentation,
        pptx_renderer.SlidePlan(kind="content", component_id="finance.variance_table", heading="Wide table", table=table),
        footer_text="Test · Deck",
    )

    cap = pptx_renderer._MAX_TABLE_ROWS_PER_SLIDE
    assert len(presentation.slides) == -(-len(rows) // cap)
    assert_shapes_in_bounds_and_non_overlapping(presentation)
    assert_font_floor(presentation)

    text = _deck_text(presentation)
    for row in rows:
        assert format_value(row.cells["value"].value, "#,##0") in text
    assert "(continued)" in text
    assert text.count("A note on the data.") == 1, "the table note should appear once, not once per page"


def test_the_divider_and_cover_are_in_bounds() -> None:
    presentation = _blank_deck()
    pptx_renderer._draw_divider(presentation, "Appendix")
    plan = pptx_renderer.PresentationPlan(
        title="A Test Deck",
        subtitle="A subtitle",
        metadata={"company": "Test Co", "author": "A. Person", "author_title": "Tester"},
        slides=(),
    )
    pptx_renderer._draw_cover(presentation, plan)
    assert_shapes_in_bounds_and_non_overlapping(presentation)
    assert_font_floor(presentation)
