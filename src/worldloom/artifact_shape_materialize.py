"""Materialize large native artifacts from :mod:`worldloom.artifact_shape` plans.

This is intentionally a *shape* materializer, not a second artifact compiler.
Business truth enters only through the fact ids named by the plan's evidence
locator. Everything else is deterministic non-evidentiary ballast used to
exercise paging, extraction, file-size and context behavior.
"""

from __future__ import annotations

import hashlib
import zipfile
from collections.abc import Mapping
from pathlib import Path
from typing import Any

from .artifact_media import as_stream, media_chunks
from .artifact_shape import ArtifactShapePlan
from .evidence_locator import DocxLocator, EvidenceRef, PptxLocator, XlsxLocator
from .models import CanonicalFact, Model


class ArtifactShapeMaterialization(Model):
    artifact_id: str
    artifact_type: str
    path: str
    file_size_bytes: int
    image_bytes: int = 0
    slides: int = 0
    pages: int = 0
    paragraphs: int = 0
    sheets: int = 0
    rows_per_sheet: int = 0
    columns_per_sheet: int = 0
    formulas: int = 0
    comments: int = 0
    evidence: EvidenceRef | None = None
    ballast_is_evidence: bool = False

    @property
    def logical_cells(self) -> int:
        return self.sheets * self.rows_per_sheet * self.columns_per_sheet


def _fact_value(fact: CanonicalFact) -> Any:
    if fact.value is not None:
        return fact.value.amount
    return fact.text_value


def _fact_text(fact: CanonicalFact) -> str:
    if fact.value is not None:
        return f"{fact.kind}: {fact.value.amount} {fact.value.unit}"
    return f"{fact.kind}: {fact.text_value}"


def _evidence_facts(
    evidence: EvidenceRef | None,
    facts: Mapping[str, CanonicalFact],
) -> tuple[CanonicalFact, ...]:
    if evidence is None:
        return ()
    missing = set(evidence.fact_ids) - set(facts)
    if missing:
        raise ValueError(f"missing evidence facts for shape materialization: {sorted(missing)}")
    return tuple(facts[fact_id] for fact_id in evidence.fact_ids)


def _evidence_text(evidence: EvidenceRef | None, facts: Mapping[str, CanonicalFact]) -> str:
    return "\n".join(_fact_text(fact) for fact in _evidence_facts(evidence, facts))


def _ballast_label(artifact_id: str, ordinal: int, *, prefix: str = "Reference appendix") -> str:
    token = hashlib.sha256(f"{artifact_id}\0{ordinal}".encode()).hexdigest()[:10]
    return f"{prefix} · {token}"


def _pad_ooxml(path: Path, target_size: int, *, seed: str) -> None:
    """Pad an OOXML ZIP to at least *target_size* without buffering the ballast.

    The extra package member is deliberately unreferenced and stored rather than
    compressed. Office readers and the Python OOXML libraries ignore it; the
    harness still has to transport/read the real file weight. A fixed ZipInfo
    timestamp keeps the padding itself byte-replayable.
    """

    if target_size <= 0 or path.stat().st_size >= target_size:
        return
    required = target_size - path.stat().st_size
    info = zipfile.ZipInfo(
        "worldloom/shape-ballast.bin",
        date_time=(1980, 1, 1, 0, 0, 0),
    )
    info.compress_type = zipfile.ZIP_STORED
    digest = hashlib.sha256(seed.encode()).digest()
    block = (digest * (4096 // len(digest) + 1))[:4096]
    with zipfile.ZipFile(path, "a", compression=zipfile.ZIP_STORED) as package:
        with package.open(info, "w") as stream:
            remaining = required
            while remaining > 0:
                chunk = block[: min(len(block), remaining)]
                stream.write(chunk)
                remaining -= len(chunk)


def _package_media_bytes(path: Path) -> int:
    with zipfile.ZipFile(path) as package:
        return sum(
            info.file_size
            for info in package.infolist()
            if info.filename.startswith("ppt/media/")
        )


def _materialize_pptx(
    plan: ArtifactShapePlan,
    out: Path,
    *,
    title: str,
    facts: Mapping[str, CanonicalFact],
) -> ArtifactShapeMaterialization:
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches, Pt
    except ImportError as error:  # pragma: no cover - optional dependency
        raise RuntimeError('PPTX shape materialization needs `pip install "worldloom[pptx]"`') from error

    slides = max(1, plan.target_slides)
    evidence = plan.evidence
    locator = evidence.locator if evidence is not None else None
    ppt_locator = locator if isinstance(locator, PptxLocator) else None
    evidence_text = _evidence_text(evidence, facts)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    note_budget = plan.target_speaker_note_slides
    hidden_budget = plan.target_hidden_slides
    chart_budget = plan.target_native_charts

    for ordinal in range(1, slides + 1):
        slide = presentation.slides.add_slide(blank)
        heading = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.7))
        heading.text_frame.text = title if ordinal == 1 else _ballast_label(plan.artifact_id, ordinal)
        heading.text_frame.paragraphs[0].runs[0].font.size = Pt(22)
        body = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.4), Inches(4.8))
        body.text_frame.text = _ballast_label(
            plan.artifact_id,
            ordinal,
            prefix="Operational reference",
        )

        is_evidence = ppt_locator is not None and ppt_locator.slide == ordinal
        if ppt_locator is not None and is_evidence and ppt_locator.element in {"text", "table", "hidden_slide"}:
            body.text_frame.text = evidence_text
        if ppt_locator is not None and is_evidence and ppt_locator.element == "hidden_slide":
            slide._element.set("show", "0")
        elif hidden_budget > 0 and ordinal > slides - hidden_budget:
            slide._element.set("show", "0")

        needs_notes = ordinal <= note_budget or (
            is_evidence and ppt_locator is not None and ppt_locator.element == "speaker_notes"
        )
        if needs_notes:
            notes = slide.notes_slide.notes_text_frame
            notes.text = (
                evidence_text
                if is_evidence and ppt_locator is not None and ppt_locator.element == "speaker_notes"
                else _ballast_label(plan.artifact_id, ordinal, prefix="Reference note")
            )

        needs_chart = ordinal <= chart_budget or (
            is_evidence and ppt_locator is not None and ppt_locator.element == "chart"
        )
        if needs_chart:
            chart_data = CategoryChartData()
            chart_data.categories = ["Reference"]
            value = 0.0
            if is_evidence and ppt_locator is not None and ppt_locator.element == "chart":
                evidence_facts = _evidence_facts(evidence, facts)
                if len(evidence_facts) != 1 or evidence_facts[0].value is None:
                    raise ValueError("chart evidence needs exactly one numeric fact")
                value = float(evidence_facts[0].value.amount)
            chart_data.add_series("actual", (value,))
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(8.6),
                Inches(4.6),
                Inches(3.2),
                Inches(1.8),
                chart_data,
            )

    if plan.target_image_bytes:
        chunks = media_chunks(f"pptx:{plan.artifact_id}:media", plan.target_image_bytes)
        for index, payload in enumerate(chunks):
            media_slide = presentation.slides[index % slides]
            media_slide.shapes.add_picture(
                as_stream(payload),
                Inches(13.0),
                Inches(7.15),
                width=Inches(0.12),
                height=Inches(0.12),
            )

    out.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(out))
    _pad_ooxml(out, plan.target_file_size_bytes, seed=f"pptx:{plan.artifact_id}")
    image_bytes = _package_media_bytes(out)
    if image_bytes < plan.target_image_bytes:
        raise ValueError(
            f"embedded media target not met: {image_bytes} < {plan.target_image_bytes}"
        )
    return ArtifactShapeMaterialization(
        artifact_id=plan.artifact_id,
        artifact_type=plan.artifact_type,
        path=str(out),
        file_size_bytes=out.stat().st_size,
        image_bytes=image_bytes,
        slides=slides,
        evidence=evidence,
        ballast_is_evidence=False,
    )


def _materialize_docx(
    plan: ArtifactShapePlan,
    out: Path,
    *,
    title: str,
    facts: Mapping[str, CanonicalFact],
) -> ArtifactShapeMaterialization:
    try:
        from docx import Document
        from docx.enum.text import WD_BREAK
    except ImportError as error:  # pragma: no cover - optional dependency
        raise RuntimeError('DOCX shape materialization needs `pip install "worldloom[docx]"`') from error

    pages = max(1, plan.target_pages)
    paragraphs = max(pages, plan.target_paragraphs or pages)
    evidence = plan.evidence
    locator = evidence.locator if evidence is not None else None
    doc_locator = locator if isinstance(locator, DocxLocator) else None
    evidence_text = _evidence_text(evidence, facts)
    document = Document()
    document.core_properties.title = title
    base, remainder = divmod(paragraphs, pages)
    global_paragraph = 0
    for page in range(1, pages + 1):
        count = base + (1 if page <= remainder else 0)
        for _ in range(count):
            global_paragraph += 1
            paragraph = document.add_paragraph()
            if doc_locator is not None and doc_locator.paragraph == global_paragraph:
                paragraph.add_run(evidence_text)
            else:
                paragraph.add_run(
                    _ballast_label(plan.artifact_id, global_paragraph, prefix="Reference note")
                )
        if page < pages:
            # An explicit break belongs to the final paragraph rather than a
            # new empty paragraph, so the native paragraph locator remains stable.
            paragraph.runs[-1].add_break(WD_BREAK.PAGE)

    out.parent.mkdir(parents=True, exist_ok=True)
    document.save(str(out))
    _pad_ooxml(out, plan.target_file_size_bytes, seed=f"docx:{plan.artifact_id}")
    return ArtifactShapeMaterialization(
        artifact_id=plan.artifact_id,
        artifact_type=plan.artifact_type,
        path=str(out),
        file_size_bytes=out.stat().st_size,
        pages=pages,
        paragraphs=paragraphs,
        evidence=evidence,
        ballast_is_evidence=False,
    )


def _spreadsheet_evidence_value(
    evidence: EvidenceRef | None,
    facts: Mapping[str, CanonicalFact],
) -> Any:
    resolved = _evidence_facts(evidence, facts)
    if len(resolved) == 1:
        return _fact_value(resolved[0])
    return " | ".join(_fact_text(fact) for fact in resolved)


def _materialize_xlsx(
    plan: ArtifactShapePlan,
    out: Path,
    *,
    title: str,
    facts: Mapping[str, CanonicalFact],
) -> ArtifactShapeMaterialization:
    try:
        from openpyxl import Workbook
        from openpyxl.cell import WriteOnlyCell
        from openpyxl.comments import Comment
    except ImportError as error:  # pragma: no cover - optional dependency
        raise RuntimeError('XLSX shape materialization needs `pip install "worldloom[xlsx]"`') from error

    sheets = max(1, plan.target_sheets)
    rows = max(1, plan.target_rows_per_sheet)
    columns = max(2 if isinstance(getattr(plan.evidence, "locator", None), XlsxLocator) else 1, plan.target_columns_per_sheet or 1)
    capacity = sheets * rows * columns
    if plan.target_formulas > capacity:
        raise ValueError("formula target exceeds logical workbook cell capacity")
    if plan.target_comments > sheets * rows:
        raise ValueError("comment target exceeds one-comment-per-row shape capacity")

    evidence = plan.evidence
    locator = evidence.locator if evidence is not None else None
    xlsx_locator = locator if isinstance(locator, XlsxLocator) else None
    evidence_value = _spreadsheet_evidence_value(evidence, facts)
    workbook = Workbook(write_only=True)
    formula_remaining = plan.target_formulas
    comment_remaining = plan.target_comments
    formula_written = 0
    comments_written = 0

    for sheet_index in range(1, sheets + 1):
        sheet = workbook.create_sheet(f"Sheet{sheet_index}")
        for row_index in range(1, rows + 1):
            width = columns if row_index == rows else 1
            if formula_remaining or comment_remaining:
                width = columns
            if xlsx_locator is not None and xlsx_locator.sheet == sheet.title and xlsx_locator.cell.endswith(str(row_index)):
                width = max(width, 2)
            values: list[Any] = [None] * width
            if row_index == 1:
                values[0] = title if sheet_index == 1 else _ballast_label(plan.artifact_id, sheet_index)
            if xlsx_locator is not None and xlsx_locator.sheet == sheet.title and xlsx_locator.cell == f"B{row_index}":
                if xlsx_locator.formula:
                    values[0] = evidence_value
                    values[1] = f"=A{row_index}"
                else:
                    values[1] = evidence_value
            for column_index in range(width):
                if formula_remaining <= 0:
                    break
                if column_index < len(values) and values[column_index] is None:
                    values[column_index] = '=""'
                    formula_remaining -= 1
                    formula_written += 1
            if comment_remaining > 0:
                cell = WriteOnlyCell(sheet, value=values[0])
                cell.comment = Comment(
                    _ballast_label(plan.artifact_id, comments_written + 1, prefix="Reference comment"),
                    "Worldloom",
                )
                values[0] = cell
                comment_remaining -= 1
                comments_written += 1
            if row_index == rows and len(values) < columns:
                values.extend([None] * (columns - len(values)))
            if row_index == rows and values[-1] is None:
                values[-1] = _ballast_label(plan.artifact_id, sheet_index, prefix="Shape boundary")
            sheet.append(values)

    out.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(out)
    _pad_ooxml(out, plan.target_file_size_bytes, seed=f"xlsx:{plan.artifact_id}")
    return ArtifactShapeMaterialization(
        artifact_id=plan.artifact_id,
        artifact_type=plan.artifact_type,
        path=str(out),
        file_size_bytes=out.stat().st_size,
        sheets=sheets,
        rows_per_sheet=rows,
        columns_per_sheet=columns,
        formulas=formula_written,
        comments=comments_written,
        evidence=evidence,
        ballast_is_evidence=False,
    )


def materialize_artifact_shape(
    plan: ArtifactShapePlan,
    out: str | Path,
    *,
    title: str,
    facts: Mapping[str, CanonicalFact] | None = None,
) -> ArtifactShapeMaterialization:
    """Create the native heavy fixture described by *plan* at *out*."""

    destination = Path(out)
    available = facts or {}
    kind = plan.artifact_type.casefold()
    if kind in {"pptx", "presentation", "gslides", "slides"}:
        return _materialize_pptx(plan, destination, title=title, facts=available)
    if kind in {"docx", "document", "gdoc", "doc"}:
        return _materialize_docx(plan, destination, title=title, facts=available)
    if kind in {"xlsx", "spreadsheet", "gsheet", "sheet"}:
        return _materialize_xlsx(plan, destination, title=title, facts=available)
    raise ValueError(f"unsupported artifact shape materializer {plan.artifact_type!r}")


__all__ = ["ArtifactShapeMaterialization", "materialize_artifact_shape"]
