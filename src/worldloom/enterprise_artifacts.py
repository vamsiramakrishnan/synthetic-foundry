"""Render real enterprise evaluation artifacts through optional libraries."""

from __future__ import annotations

from collections.abc import Iterable
from pathlib import Path
from typing import TYPE_CHECKING, Any

from .enterprise_corpus import EnterpriseCorpus, QueryFixture
from .enterprise_queries import PlannedEnterpriseQuery
from .models import Model

if TYPE_CHECKING:
    from .connector_data import ConnectorRecord


class RenderedEvalArtifact(Model):
    query_id: str
    format: str
    path: str
    source_record_ids: tuple[str, ...]


def _source_records(
    corpus: EnterpriseCorpus, fixture: QueryFixture
) -> tuple[ConnectorRecord, ...]:
    identifiers = {
        identifier
        for values in fixture.input_record_ids.values()
        for identifier in values
    }
    return tuple(
        record
        for record in corpus.connector_data.records
        if record.id in identifiers
    )


def _rows(records: Iterable[ConnectorRecord]) -> list[dict[str, Any]]:
    rows = []
    for record in records:
        rows.append(
            {
                "connector": record.connector,
                "entity": record.entity,
                "external_id": record.external_id,
                "title": record.title,
                "facts": len(record.fact_ids),
                "events": len(record.event_ids),
            }
        )
    return rows


def _render_xlsx(
    query: PlannedEnterpriseQuery, records: tuple[ConnectorRecord, ...], path: Path
) -> None:
    try:
        import xlsxwriter  # type: ignore[import-untyped]
    except ImportError as error:
        raise RuntimeError("XLSX rendering requires worldloom[xlsx]") from error
    workbook = xlsxwriter.Workbook(path)
    title = workbook.add_format({"bold": True, "font_size": 16})
    header = workbook.add_format({"bold": True, "bg_color": "#D9EAF7"})
    summary = workbook.add_worksheet("Summary")
    summary.write("A1", query.workflow.replace("_", " ").title(), title)
    summary.write("A3", "Customer request", header)
    summary.write("A4", query.query)
    rows = _rows(records)
    detail = workbook.add_worksheet("Detail")
    columns = ("connector", "entity", "external_id", "title", "facts", "events")
    for column, name in enumerate(columns):
        detail.write(0, column, name, header)
    for row_index, row in enumerate(rows, start=1):
        for column, name in enumerate(columns):
            detail.write(row_index, column, row[name])
    detail.add_table(0, 0, max(len(rows), 1), len(columns) - 1, {"name": "Evidence", "columns": [{"header": item} for item in columns]})
    counts: dict[str, int] = {}
    for row in rows:
        counts[row["connector"]] = counts.get(row["connector"], 0) + 1
    chart_data = workbook.add_worksheet("Chart Data")
    chart_data.write_row(0, 0, ("Connector", "Records"), header)
    for index, (connector, count) in enumerate(sorted(counts.items()), start=1):
        chart_data.write_row(index, 0, (connector, count))
    chart = workbook.add_chart({"type": "column"})
    if counts:
        chart.add_series({"name": "Records", "categories": ["Chart Data", 1, 0, len(counts), 0], "values": ["Chart Data", 1, 1, len(counts), 1]})
    chart.set_title({"name": "Evidence records by connector"})
    summary.insert_chart("A7", chart)
    provenance = workbook.add_worksheet("Provenance")
    provenance.write_row(0, 0, ("Record ID", "Fact IDs", "Event IDs"), header)
    for index, record in enumerate(records, start=1):
        provenance.write_row(index, 0, (record.id, ",".join(record.fact_ids), ",".join(record.event_ids)))
    workbook.close()


def _render_docx(
    query: PlannedEnterpriseQuery, records: tuple[ConnectorRecord, ...], path: Path
) -> None:
    try:
        from docx import Document
    except ImportError as error:
        raise RuntimeError("DOCX rendering requires worldloom[docx]") from error
    document = Document()
    document.add_heading(query.workflow.replace("_", " ").title(), 0)
    document.add_heading("Executive summary", level=1)
    document.add_paragraph(query.query)
    document.add_heading("Evidence", level=1)
    table = document.add_table(rows=1, cols=4)
    for cell, value in zip(table.rows[0].cells, ("Connector", "Entity", "Record", "Title"), strict=True):
        cell.text = value
    for record in records:
        cells = table.add_row().cells
        for cell, value in zip(cells, (record.connector, record.entity, record.external_id, record.title), strict=True):
            cell.text = value
    document.add_heading("Sources", level=1)
    for record in records:
        document.add_paragraph(f"{record.connector}:{record.external_id} — facts: {', '.join(record.fact_ids) or 'none'}", style="List Bullet")
    document.save(str(path))


def _render_pptx(
    query: PlannedEnterpriseQuery, records: tuple[ConnectorRecord, ...], path: Path
) -> None:
    try:
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
    except ImportError as error:
        raise RuntimeError("PPTX rendering requires worldloom[pptx]") from error
    deck = Presentation()
    title = deck.slides.add_slide(deck.slide_layouts[0])
    title.shapes.title.text = query.workflow.replace("_", " ").title()
    title.placeholders[1].text = query.query
    summary = deck.slides.add_slide(deck.slide_layouts[1])
    summary.shapes.title.text = "Executive summary"
    summary.placeholders[1].text = f"Synthesized from {len(records)} grounded connector records."
    counts: dict[str, int] = {}
    for record in records:
        counts[record.connector] = counts.get(record.connector, 0) + 1
    chart_slide = deck.slides.add_slide(deck.slide_layouts[5])
    chart_slide.shapes.title.text = "Evidence by connector"
    data = ChartData()
    data.categories = list(sorted(counts))
    data.add_series("Records", [counts[item] for item in sorted(counts)])
    chart_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(8), Inches(4.5), data)
    sources = deck.slides.add_slide(deck.slide_layouts[1])
    sources.shapes.title.text = "Sources"
    sources.placeholders[1].text = "\n".join(f"{record.connector}:{record.external_id} — {record.title}" for record in records[:20])
    deck.save(str(path))


def _render_pdf(
    query: PlannedEnterpriseQuery, records: tuple[ConnectorRecord, ...], path: Path
) -> None:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table
    except ImportError as error:
        raise RuntimeError("PDF rendering requires worldloom[pdf]") from error
    styles = getSampleStyleSheet()
    story: list[Any] = [Paragraph(query.workflow.replace("_", " ").title(), styles["Title"]), Spacer(1, 12), Paragraph(query.query, styles["BodyText"]), Spacer(1, 12)]
    data = [["Connector", "Entity", "Record", "Title"]]
    data.extend([[record.connector, record.entity, record.external_id, record.title] for record in records])
    story.append(Table(data, repeatRows=1))
    SimpleDocTemplate(str(path), pagesize=A4).build(story)


def render_eval_artifact(
    query: PlannedEnterpriseQuery,
    fixture: QueryFixture,
    corpus: EnterpriseCorpus,
    directory: Path,
) -> RenderedEvalArtifact | None:
    requirement = query.generation.artifact
    if requirement is None or requirement.format not in {"xlsx", "docx", "pptx", "pdf"}:
        return None
    directory.mkdir(parents=True, exist_ok=True)
    records = _source_records(corpus, fixture)
    path = directory / f"{query.id}.{requirement.format}"
    renderer = {"xlsx": _render_xlsx, "docx": _render_docx, "pptx": _render_pptx, "pdf": _render_pdf}[requirement.format]
    renderer(query, records, path)
    return RenderedEvalArtifact(query_id=query.id, format=requirement.format, path=str(path), source_record_ids=tuple(record.id for record in records))


def render_corpus_artifacts(
    corpus: EnterpriseCorpus, directory: Path, *, limit: int | None = None
) -> tuple[RenderedEvalArtifact, ...]:
    fixtures = {fixture.query_id: fixture for fixture in corpus.fixtures}
    rendered: list[RenderedEvalArtifact] = []
    for query in corpus.queries:
        if limit is not None and len(rendered) >= limit:
            break
        artifact = render_eval_artifact(query, fixtures[query.id], corpus, directory)
        if artifact is not None:
            rendered.append(artifact)
    return tuple(rendered)
